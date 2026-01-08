"""
Generate BSc Thesis Defense Presentation
Creates PowerPoint slides for 20-minute defense talk
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """Title slide with gradient background"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Background gradient (teal to cyan)
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 45.0
    fill.gradient_stops[0].color.rgb = RGBColor(13, 148, 136)  # Teal-600
    fill.gradient_stops[1].color.rgb = RGBColor(6, 182, 212)  # Cyan-500
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(48)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.name = "Segoe UI"
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(26)
    subtitle_frame.paragraphs[0].font.name = "Segoe UI Light"
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(204, 251, 241)  # Teal-100
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Author & date
    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(0.8))
    info_frame = info_box.text_frame
    info_frame.text = "BSc Thesis Defense | December 2024"
    info_frame.paragraphs[0].font.size = Pt(18)
    info_frame.paragraphs[0].font.name = "Segoe UI"
    info_frame.paragraphs[0].font.color.rgb = RGBColor(165, 243, 252)  # Cyan-200
    info_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullet_points):
    """Content slide with modern card design"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Background (soft cream)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 251, 235)  # Amber-50
    
    # Title bar (coral/orange)
    header = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.5), Inches(0.4), Inches(9), Inches(0.9)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(251, 146, 60)  # Orange-400
    header.line.fill.background()
    header.shadow.inherit = False
    
    # Title on header
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(8.6), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(30)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.name = "Segoe UI Semibold"
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Content card
    card = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.7), Inches(1.6), Inches(8.6), Inches(5.3)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    card.line.color.rgb = RGBColor(251, 191, 36)  # Amber-400
    card.line.width = Pt(2)
    
    # Content text
    content_box = slide.shapes.add_textbox(Inches(1.2), Inches(2), Inches(7.6), Inches(4.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, bullet in enumerate(bullet_points):
        p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
        p.text = "▪ " + bullet  # Small black square bullet
        p.font.size = Pt(19)
        p.font.name = "Calibri"
        p.font.color.rgb = RGBColor(30, 41, 59)  # Slate-800
        p.space_after = Pt(14)
        p.line_spacing = 1.3
    
    return slide

def add_image_slide(prs, title, image_path, caption=""):
    """Slide with image and modern card design"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Background (soft cream)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 251, 235)  # Amber-50
    
    # Title bar (teal)
    header = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.5), Inches(0.4), Inches(9), Inches(0.75)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(20, 184, 166)  # Teal-500
    header.line.fill.background()
    
    # Title on header
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(8.6), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(26)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.name = "Segoe UI Semibold"
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Image card
    card = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.6), Inches(1.4), Inches(8.8), Inches(5.4)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    card.line.color.rgb = RGBColor(45, 212, 191)  # Teal-400
    card.line.width = Pt(2)
    
    # Image - use absolute path
    img_path = Path(image_path)
    if not img_path.is_absolute():
        img_path = Path("thesis") / img_path
    if img_path.exists():
        pic = slide.shapes.add_picture(str(img_path.absolute()), Inches(0.8), Inches(1.6), width=Inches(8.4))
    else:
        # Add placeholder text if image not found
        error_box = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(6), Inches(1))
        error_frame = error_box.text_frame
        error_frame.text = f"[Image: {image_path}]"
        error_frame.paragraphs[0].font.size = Pt(18)
        error_frame.paragraphs[0].font.name = "Calibri"
        error_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        error_frame.paragraphs[0].font.color.rgb = RGBColor(100, 116, 139)  # Slate-500
    
    # Caption
    if caption:
        caption_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.9), Inches(8.4), Inches(0.5))
        caption_frame = caption_box.text_frame
        caption_frame.text = caption
        caption_frame.paragraphs[0].font.size = Pt(15)
        caption_frame.paragraphs[0].font.name = "Calibri"
        caption_frame.paragraphs[0].font.color.rgb = RGBColor(71, 85, 105)  # Slate-600
        caption_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

# ============================================================================
# Slide 1: Title
# ============================================================================
add_title_slide(prs, 
    "Fair and Robust Machine Learning",
    "Achieving Perfect Fairness Through Iterative Adaptive Sample Weighting"
)

# ============================================================================
# Slide 2: Motivation - The Fairness Problem
# ============================================================================
add_content_slide(prs, "The Fairness Problem", [
    "ML systems make high-stakes decisions: loans, hiring, criminal justice",
    "Real-world bias examples:",
    "  • COMPAS recidivism: 45% false positive rate for African-Americans vs 23% for Caucasians",
    "  • Amazon hiring: AI rejected female candidates due to historical bias",
    "Challenge: How to achieve fairness without sacrificing accuracy?"
])

# ============================================================================
# Slide 3: Research Questions
# ============================================================================
add_content_slide(prs, "Research Questions", [
    "RQ1: Can iterative adaptive weighting achieve perfect fairness?",
    "RQ2: What are the trade-offs (accuracy, calibration)?",
    "RQ3: Is it computationally feasible for production?",
    "RQ4: How does the mechanism work?"
])

# ============================================================================
# Slide 4: Our Approach
# ============================================================================
add_content_slide(prs, "Our Approach: Iterative Adaptive Weighting", [
    "Weight formula: w_i = (c_i × r_i + ε)^(1/T)",
    "  • c_i = confidence (distance from decision boundary)",
    "  • r_i = correctness (1 if correct, 0 if wrong)",
    "  • T = temperature (controls concentration)",
    "Counterintuitive: Upweights confident correct predictions",
    "Iterative: Recompute weights after each training iteration"
])

# ============================================================================
# Slide 5: Algorithm
# ============================================================================
add_content_slide(prs, "Iterative Training Algorithm", [
    "1. Initialize: All weights w_i = 1",
    "2. Train model with current weights",
    "3. Compute confidence c_i and correctness r_i for each sample",
    "4. Update weights: w_i = (c_i × r_i + ε)^(1/T)",
    "5. Repeat until fairness threshold achieved (EO < 0.01)",
    "Typically converges in 4-10 iterations, <2 seconds training"
])

# ============================================================================
# Slide 6: Key Result - Perfect Fairness
# ============================================================================
add_content_slide(prs, "Key Result: Perfect Fairness Achieved", [
    "German Credit Dataset:",
    "  ✓ Equalized Odds = 0.000 (perfect fairness!)",
    "  ✓ Demographic Parity = 0.000",
    "  ✓ First reported perfect fairness on real data (in-processing method)",
    "",
    "Adult Income Dataset:",
    "  ✓ 68.7% reduction in EO violations (0.163 → 0.051)",
    "  ✓ Competitive with post-processing methods"
])

# ============================================================================
# Slide 7: Fairness Comparison Figure
# ============================================================================
add_image_slide(prs, 
    "Fairness Performance Across Datasets",
    "figures/fairness_comparison.png",
    "Our method (blue) achieves perfect fairness on German, substantial improvement on Adult"
)

# ============================================================================
# Slide 8: The Trade-off - Calibration Degradation
# ============================================================================
add_content_slide(prs, "Trade-off: Calibration Degradation", [
    "Perfect fairness comes at a cost:",
    "  • German: +388% ECE increase (0.089 → 0.434)",
    "  • Adult: +756% ECE increase (0.052 → 0.445)",
    "",
    "Fundamental tension between fairness and calibration",
    "First systematic quantification for in-processing methods",
    "Acceptable for decision-based apps, problematic for probability-based"
])

# ============================================================================
# Slide 9: Reliability Diagrams
# ============================================================================
add_image_slide(prs,
    "Calibration Degradation Visualized",
    "figures/reliability_diagrams.png",
    "Left: Baseline well-calibrated. Right: Our method shows overconfidence at high probabilities"
)

# ============================================================================
# Slide 10: Accuracy Trade-off (Minimal)
# ============================================================================
add_content_slide(prs, "Accuracy Trade-off: Minimal Impact", [
    "Accuracy degradation is small:",
    "  • German: -1.8% (0.724 → 0.706)",
    "  • Adult: -0.9% (0.851 → 0.842)",
    "  • COMPAS: -0.3% (0.673 → 0.670)",
    "",
    "Contradicts belief that perfect fairness requires large accuracy loss",
    "Fairness-accuracy trade-off is mild"
])

# ============================================================================
# Slide 11: How It Works - Mechanism
# ============================================================================
add_content_slide(prs, "How Does It Work? Mechanism Insights", [
    "Key insight: Exploits confidence asymmetry between groups",
    "  • Disadvantaged groups have lower confidence even for correct predictions",
    "  • Weight formula amplifies this via exponent 1/T",
    "",
    "Mechanism:",
    "  1. Upweight confident correct predictions from disadvantaged group",
    "  2. Model 'pays more attention' to these samples",
    "  3. Over iterations, TPR/FPR gaps close → fairness achieved"
])

# ============================================================================
# Slide 12: Weight Distribution Evolution
# ============================================================================
add_image_slide(prs,
    "Weight Evolution Across Iterations",
    "figures/weight_distribution.png",
    "Iteration 1: uniform. Iteration 5: concentrated on confident correct predictions"
)

# ============================================================================
# Slide 13: Computational Efficiency
# ============================================================================
add_content_slide(prs, "Computational Efficiency", [
    "Training time: <2 seconds for 32K samples (Adult dataset)",
    "Iterations to convergence: 4-10 (German converges in 4)",
    "Scalability: O(n) linear time complexity",
    "",
    "Zero inference overhead:",
    "  • No test-time modifications needed",
    "  • Standard model serving infrastructure works unchanged",
    "  • Simplifies production deployment"
])

# ============================================================================
# Slide 14: Dataset Dependency
# ============================================================================
add_content_slide(prs, "Limitation: Dataset-Dependent Effectiveness", [
    "Results vary across datasets:",
    "  ✓ German: Perfect fairness (EO = 0.000)",
    "  ✓ Adult: Substantial improvement (68.7% reduction)",
    "  ✗ COMPAS: Limited success (14.6% reduction)",
    "",
    "Hypothesis: Effectiveness depends on:",
    "  • Group imbalance structure",
    "  • Base rate differences between groups",
    "  • Feature space complexity"
])

# ============================================================================
# Slide 15: Limitations Summary
# ============================================================================
add_content_slide(prs, "Limitations", [
    "1. Severe calibration degradation (+388-756% ECE)",
    "   → Not suitable for probability-based applications",
    "",
    "2. Dataset-dependent effectiveness",
    "   → Pilot testing required for each dataset",
    "",
    "3. Binary classification with single sensitive attribute",
    "   → Extension to multi-class, intersectional fairness needed"
])

# ============================================================================
# Slide 16: Contributions Summary
# ============================================================================
add_content_slide(prs, "Contributions", [
    "1. Perfect fairness on real-world data (first for in-processing)",
    "2. Fairness-calibration trade-off quantified (+388-756% ECE)",
    "3. Novel interpretable mechanism (confidence × correctness)",
    "4. Zero inference overhead (deployment simplicity)",
    "5. Computational efficiency (<2s training, 4-10 iterations)",
    "6. Open-source implementation (reproducible)"
])

# ============================================================================
# Slide 17: Practical Implications
# ============================================================================
add_content_slide(prs, "Practical Implications", [
    "For practitioners:",
    "  • Simple method (10 lines of code) can achieve SOTA fairness",
    "  • Zero production overhead makes deployment trivial",
    "  • Choose based on application: decision-based ✓, probability-based ✗",
    "",
    "For fairness research:",
    "  • Fairness-calibration trade-off is fundamental",
    "  • Perfect fairness is achievable (not just asymptotic)",
    "  • Dataset characteristics matter for method effectiveness"
])

# ============================================================================
# Slide 18: Future Work
# ============================================================================
add_content_slide(prs, "Future Work", [
    "1. Integrated recalibration: Restore calibration post-hoc",
    "2. Extension to neural networks: Deep learning compatibility",
    "3. Intersectional fairness: Multiple sensitive attributes (race × gender)",
    "4. Theoretical analysis: Convergence guarantees, optimality",
    "5. Real-world deployment study: Production validation"
])

# ============================================================================
# Slide 19: Conclusion
# ============================================================================
add_content_slide(prs, "Conclusion", [
    "Iterative adaptive weighting achieves perfect fairness:",
    "  ✓ EO = 0.000 on German Credit (first for in-processing)",
    "  ✓ Minimal accuracy loss (-0.9% to -1.8%)",
    "  ✓ Zero inference overhead",
    "",
    "But at a cost:",
    "  ✗ Severe calibration degradation (+388-756% ECE)",
    "  ✗ Dataset-dependent effectiveness",
    "",
    "Trade-offs are measurable, navigable, and application-dependent"
])

# ============================================================================
# Slide 20: Thank You / Questions
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Thank you text
thank_you_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
thank_you_frame = thank_you_box.text_frame
thank_you_frame.text = "Thank You!"
thank_you_frame.paragraphs[0].font.size = Pt(60)
thank_you_frame.paragraphs[0].font.bold = True
thank_you_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Questions text
questions_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1))
questions_frame = questions_box.text_frame
questions_frame.text = "Questions?"
questions_frame.paragraphs[0].font.size = Pt(40)
questions_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Contact info
contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
contact_frame = contact_box.text_frame
contact_frame.text = "Code & Thesis: github.com/[your-username]/fair-robust-thesis"
contact_frame.paragraphs[0].font.size = Pt(14)
contact_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ============================================================================
# Save presentation
# ============================================================================
output_path = Path("thesis/defense_presentation.pptx")
prs.save(output_path)

print("=" * 60)
print("THESIS DEFENSE PRESENTATION CREATED")
print("=" * 60)
print(f"\nFile: {output_path.absolute()}")
print(f"Slides: 20 (Title + 18 content + Q&A)")
print(f"Duration: ~20 minutes (1 min/slide)")
print("\nSlide Breakdown:")
print("  1. Title")
print("  2-3. Motivation & Research Questions (3 min)")
print("  4-5. Method & Algorithm (3 min)")
print("  6-9. Results: Fairness & Calibration (6 min)")
print("  10-12. Mechanism & Efficiency (4 min)")
print("  13-15. Limitations (3 min)")
print("  16-19. Contributions & Future Work (3 min)")
print("  20. Q&A")
print("\n" + "=" * 60)
print("\nNext steps:")
print("1. Open defense_presentation.pptx in PowerPoint")
print("2. Review slides, adjust formatting as needed")
print("3. Practice 20-minute talk with timer")
print("4. Prepare answers to anticipated questions")
print("=" * 60)
